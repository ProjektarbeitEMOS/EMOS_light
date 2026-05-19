"""Erzeugt eine PowerPoint zur Funktionsweise des Backend-Optimierers.

Aufbau: Zielfunktion -> Komponenten-Beiträge -> Solver -> Ausblick.
Formeln werden via matplotlib mathtext gerendert und als Bild eingebettet.
"""

import io
import os
import sys

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# ----------------------------------------------------------------------
# Farben (konsistent mit den PDFs)
# ----------------------------------------------------------------------

C_DARK_BLUE = RGBColor(0x0B, 0x3D, 0x91)
C_MID_BLUE = RGBColor(0x14, 0x3F, 0x7A)
C_GREY = RGBColor(0x55, 0x55, 0x55)
C_DARK_GREY = RGBColor(0x33, 0x33, 0x33)
C_LIGHT_GREY = RGBColor(0xF5, 0xF7, 0xFC)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN = RGBColor(0x0A, 0x6E, 0x2E)
C_AMBER = RGBColor(0xB8, 0x6E, 0x00)
C_LIGHT_GREEN_BG = RGBColor(0xEA, 0xF6, 0xED)
C_LIGHT_AMBER_BG = RGBColor(0xFB, 0xF3, 0xE1)
C_LIGHT_BLUE_BG = RGBColor(0xEE, 0xF2, 0xFA)

# 16:9 Slide
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ----------------------------------------------------------------------
# Equation rendering
# ----------------------------------------------------------------------

def render_equation(latex: str, fontsize: int = 22, dpi: int = 220) -> bytes:
    """Rendert LaTeX-Mathstring als transparentes PNG (Bytes)."""
    fig = plt.figure(figsize=(0.01, 0.01))
    fig.patch.set_alpha(0.0)
    fig.text(0, 0, f"${latex}$", fontsize=fontsize, color="black")
    buf = io.BytesIO()
    fig.savefig(
        buf, format="png", dpi=dpi,
        bbox_inches="tight", pad_inches=0.05, transparent=True,
    )
    plt.close(fig)
    buf.seek(0)
    return buf.getvalue()


# ----------------------------------------------------------------------
# Slide-Bausteine
# ----------------------------------------------------------------------

def add_blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # leere Folie
    return prs.slides.add_slide(blank_layout)


def add_title_bar(slide, title: str, subtitle: str = ""):
    """Standard-Titelleiste oben auf jeder Slide."""
    # Hintergrundbalken
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_W, Inches(0.9),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_DARK_BLUE
    bar.line.fill.background()
    bar.shadow.inherit = False

    # Titel
    tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.1), SLIDE_W - Inches(0.8), Inches(0.7),
    )
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = RGBColor(0xCC, 0xD7, 0xEC)
        r2.font.name = "Calibri"


def add_footer(slide, page_num: int):
    tb = slide.shapes.add_textbox(
        Inches(0.4), SLIDE_H - Inches(0.45),
        SLIDE_W - Inches(0.8), Inches(0.3),
    )
    tf = tb.text_frame
    tf.margin_left = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = f"EMOS Light  —  Backend-Optimierer"
    r.font.size = Pt(9)
    r.font.color.rgb = C_GREY
    r.font.name = "Calibri"
    # Seitenzahl rechts
    tb2 = slide.shapes.add_textbox(
        Inches(0.4), SLIDE_H - Inches(0.45),
        SLIDE_W - Inches(0.8), Inches(0.3),
    )
    tf2 = tb2.text_frame
    tf2.margin_left = Pt(0)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = str(page_num)
    r2.font.size = Pt(9)
    r2.font.color.rgb = C_GREY
    r2.font.name = "Calibri"


def add_bullets(slide, items, left=Inches(0.6), top=Inches(1.4),
                width=None, height=None, fontsize=18, line_spacing=1.15):
    if width is None:
        width = SLIDE_W - Inches(1.2)
    if height is None:
        height = Inches(5.5)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)

    for i, item in enumerate(items):
        if isinstance(item, str):
            text = item
            indent = 0
            bold_lead = False
        else:
            text = item.get("text", "")
            indent = item.get("indent", 0)
            bold_lead = item.get("bold_lead", False)

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = indent
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(6)

        # Bullet-Punkt manuell vor dem Text
        bullet = "•" if indent == 0 else "–"
        r0 = p.add_run()
        r0.text = f"{bullet}  "
        r0.font.size = Pt(fontsize)
        r0.font.color.rgb = C_DARK_BLUE
        r0.font.bold = True
        r0.font.name = "Calibri"

        # Text mit eventuellem Lead-Bold
        if bold_lead and ":" in text:
            head, rest = text.split(":", 1)
            r1 = p.add_run()
            r1.text = f"{head}:"
            r1.font.size = Pt(fontsize)
            r1.font.bold = True
            r1.font.color.rgb = C_DARK_GREY
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(fontsize)
            r2.font.color.rgb = C_DARK_GREY
            r2.font.name = "Calibri"
        else:
            r1 = p.add_run()
            r1.text = text
            r1.font.size = Pt(fontsize)
            r1.font.color.rgb = C_DARK_GREY
            r1.font.name = "Calibri"


def add_equation_image(slide, latex: str, left, top, target_height_in: float,
                        fontsize: int = 22):
    """Rendert Formel und platziert sie an gegebener Position."""
    png_bytes = render_equation(latex, fontsize=fontsize)
    stream = io.BytesIO(png_bytes)
    pic = slide.shapes.add_picture(stream, left, top)
    # Auf Zielhöhe skalieren
    target_h = Inches(target_height_in)
    if pic.height != target_h:
        ratio = target_h / pic.height
        pic.height = int(pic.height * ratio)
        pic.width = int(pic.width * ratio)
    return pic


def add_callout_box(slide, text: str, left, top, width, height,
                     bg_color=C_LIGHT_BLUE_BG, border_color=C_MID_BLUE,
                     text_color=C_DARK_GREY, fontsize=14, bold_first_word=True):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.color.rgb = border_color
    box.line.width = Pt(0.8)
    box.shadow.inherit = False

    tf = box.text_frame
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fontsize)
    r.font.color.rgb = text_color
    r.font.name = "Calibri"


def add_section_title(slide, big_title: str, kicker: str = ""):
    """Großer Section-Übergangs-Titel."""
    # Vollflächiger blauer Hintergrund
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK_BLUE
    bg.line.fill.background()

    if kicker:
        kb = slide.shapes.add_textbox(
            Inches(0.8), Inches(2.5), SLIDE_W - Inches(1.6), Inches(0.6),
        )
        ktf = kb.text_frame
        kp = ktf.paragraphs[0]
        kp.alignment = PP_ALIGN.LEFT
        kr = kp.add_run()
        kr.text = kicker
        kr.font.size = Pt(20)
        kr.font.color.rgb = RGBColor(0xCC, 0xD7, 0xEC)
        kr.font.name = "Calibri"

    tb = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.0), SLIDE_W - Inches(1.6), Inches(2.0),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = big_title
    r.font.size = Pt(54)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Calibri"


# ----------------------------------------------------------------------
# Die einzelnen Slides
# ----------------------------------------------------------------------

def slide_cover(prs):
    slide = add_blank_slide(prs)

    # Voller Hintergrund
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK_BLUE
    bg.line.fill.background()

    # Title
    tb = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.4), SLIDE_W - Inches(2.0), Inches(1.5),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "EMOS Light"
    r.font.size = Pt(72)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Calibri"

    # Subtitle
    tb2 = slide.shapes.add_textbox(
        Inches(1.0), Inches(3.8), SLIDE_W - Inches(2.0), Inches(1.0),
    )
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "Der Backend-Optimierer"
    r2.font.size = Pt(36)
    r2.font.color.rgb = RGBColor(0xCC, 0xD7, 0xEC)
    r2.font.name = "Calibri"

    # Untertitel
    tb3 = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.7), SLIDE_W - Inches(2.0), Inches(1.0),
    )
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    r3 = p3.add_run()
    r3.text = "Von der Zielfunktion bis zur Lösung"
    r3.font.size = Pt(20)
    r3.font.color.rgb = RGBColor(0xCC, 0xD7, 0xEC)
    r3.font.italic = True
    r3.font.name = "Calibri"

    # Footer Meta
    tb4 = slide.shapes.add_textbox(
        Inches(1.0), Inches(6.6), SLIDE_W - Inches(2.0), Inches(0.5),
    )
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    p4.alignment = PP_ALIGN.LEFT
    r4 = p4.add_run()
    r4.text = "Projektarbeit EMOS Light   |   Mai 2026"
    r4.font.size = Pt(14)
    r4.font.color.rgb = RGBColor(0xAA, 0xB8, 0xD2)
    r4.font.name = "Calibri"


def slide_agenda(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Agenda",
                  "Was wir uns in den nächsten Minuten anschauen")
    items = [
        {"text": "Was tut der Optimierer? Aufgabe in einem Satz", "bold_lead": False},
        {"text": "Die Zielfunktion – vier Kostenterme, die der Solver minimiert", "bold_lead": False},
        {"text": "Wer trägt was bei? Beiträge der einzelnen Komponenten", "bold_lead": False},
        {"text": "Die Knotenbilanz – die zentrale Kopplung", "bold_lead": False},
        {"text": "Wie löst der Solver das? LP-Relaxation, Branch-and-Bound, Schnittebenen", "bold_lead": False},
        {"text": "Modellgröße und Laufzeit bei EMOS Light", "bold_lead": False},
        {"text": "Was steht – und was noch ansteht", "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=20, line_spacing=1.4)
    add_footer(slide, page)


def slide_what_does_it_do(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Was tut der Optimierer?",
                  "Eine Aufgabe, viele Stellschrauben")

    items = [
        "Bestimmt die Steuerung aller regelbaren Energiekomponenten "
        "über den Planungshorizont (typisch 24 h, in 15-min-Schritten)",
        "Ziel: Gesamtkosten minimieren – Netzbezug + Komfortverletzung + "
        "Batterie-Verschleiß minus Einspeise-Erlös",
        "Unter Einhaltung aller physikalischen und komfortbezogenen "
        "Nebenbedingungen",
        "Mathematische Klasse: Mixed-Integer Linear Program (MILP)",
        "Solver: HiGHS (Open Source, MIT-Lizenz, schnell)",
    ]
    add_bullets(slide, items, fontsize=20, line_spacing=1.4,
                 top=Inches(1.3))

    # Visualisierung: Inputs -> Black Box -> Output
    box_y = Inches(5.4)
    # Inputs
    add_callout_box(
        slide, "Eingaben\nStrompreise, Wetter, Lastprofil, "
        "Komponenten-Setup",
        Inches(0.6), box_y, Inches(3.5), Inches(1.3),
        bg_color=C_LIGHT_GREY, border_color=C_GREY, fontsize=12,
    )
    # Optimizer
    add_callout_box(
        slide, "MILP-Optimierer\nmin Σ Kosten unter Constraints",
        Inches(4.5), box_y, Inches(4.3), Inches(1.3),
        bg_color=C_DARK_BLUE, border_color=C_DARK_BLUE,
        text_color=C_WHITE, fontsize=13,
    )
    # Output
    add_callout_box(
        slide, "Ausgabe\n96 Steuerwerte je Komponente "
        "+ Kosten-KPIs",
        Inches(9.2), box_y, Inches(3.5), Inches(1.3),
        bg_color=C_LIGHT_GREEN_BG, border_color=C_GREEN, fontsize=12,
    )

    add_footer(slide, page)


def slide_section_obj_func(prs, page):
    slide = add_blank_slide(prs)
    add_section_title(slide, "Die Zielfunktion",
                      kicker="Teil 1")


def slide_objective_overview(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Zielfunktion im Überblick",
                  "Vier Kostenterme – einer mit negativem Vorzeichen")

    # Hauptformel
    add_equation_image(
        slide,
        r"\min_{x \in X} \; K^{\text{Netz}} + K^{\text{Slack}} + K^{\text{age}}",
        Inches(2.5), Inches(1.4), target_height_in=0.9, fontsize=24,
    )

    items = [
        {"text": "Netzkosten: Bezug kostet, Einspeisung erlöst (negativer Term)",
         "bold_lead": True},
        {"text": "Komfort-Strafkosten: hohe Strafe für unerfüllten Wärmebedarf",
         "bold_lead": True},
        {"text": "Batterie-Alterung: bremst übermäßiges Cycling",
         "bold_lead": True},
    ]
    add_bullets(slide, items, fontsize=20, line_spacing=1.5,
                 top=Inches(3.3))

    # Hinweis
    add_callout_box(
        slide,
        "Alles in ct, am Ende durch 100 für EUR. Die Zielfunktion ist "
        "linear in den Entscheidungsvariablen – darum bleibt das Problem "
        "lösbar als MILP.",
        Inches(0.6), Inches(6.1), SLIDE_W - Inches(1.2), Inches(0.85),
        bg_color=C_LIGHT_AMBER_BG, border_color=C_AMBER,
        text_color=C_DARK_GREY, fontsize=14,
    )

    add_footer(slide, page)


def slide_term_grid(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Term 1 + 2:  Netzkosten und Einspeise-Erlös",
                  "Das Hauptmotiv des Solvers")

    add_equation_image(
        slide,
        r"K^{\text{Netz}} = \sum_{t \in T} \Delta t \cdot "
        r"\left[ \pi_t \cdot P^{\text{buy}}_t - \pi^{\text{feed}} \cdot P^{\text{sell}}_t \right]",
        Inches(2.0), Inches(1.4), target_height_in=0.9, fontsize=24,
    )

    items = [
        "π_t: dynamischer Endkundenpreis je Zeitschritt "
        "(Day-Ahead-Börse + Aufschläge + Steuern + MwSt)",
        "π^feed: konstante Einspeisevergütung "
        "(typisch 7–8 ct/kWh, EEG-Festvergütung)",
        "Δt: Zeitschrittlänge in Stunden (z.B. 0.25 bei 15-min-Schritten)",
        "Bezug kostet, Einspeisung bringt Erlös – der Solver "
        "balanciert beide gegeneinander",
    ]
    add_bullets(slide, items, fontsize=18, line_spacing=1.35,
                 top=Inches(3.0))

    add_callout_box(
        slide,
        "Praktischer Effekt: lieber PV selbst nutzen (oder in die Batterie/WP "
        "speichern), als für ~8 ct einspeisen, wenn später für 30 ct bezogen "
        "werden müsste.",
        Inches(0.6), Inches(6.1), SLIDE_W - Inches(1.2), Inches(0.85),
        fontsize=14,
    )

    add_footer(slide, page)


def slide_term_slack(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Term 3:  Komfort-Strafkosten (Slack)",
                  "Sicherheitsventil, das fast nie auslöst")

    add_equation_image(
        slide,
        r"K^{\text{Slack}} = c^{\text{slack}} \sum_{t \in T} \Delta t \cdot "
        r"\left( s^{\text{Heiz}}_t + s^{\text{WW}}_t \right)",
        Inches(2.5), Inches(1.4), target_height_in=0.9, fontsize=24,
    )

    items = [
        "s^Heiz_t, s^WW_t: Slack-Variablen ≥ 0 für nicht-gedeckten "
        "Heiz- bzw. Warmwasserbedarf",
        "c^slack = 500 ct/kWh – bewusst sehr hoch gewählt",
        "Damit verletzt der Solver Komfort nur, wenn das Problem sonst "
        "infeasible wäre",
        "Slack > 0 im Ergebnis ist immer ein Warnsignal: zu wenig Kapazität, "
        "zu strenge Constraints oder unrealistische Eingaben",
    ]
    add_bullets(slide, items, fontsize=18, line_spacing=1.35,
                 top=Inches(3.0))

    add_callout_box(
        slide,
        "Warum 500 ct/kWh?  Selbst im teuersten Strom-Szenario "
        "(~50 ct/kWh) lohnt es sich für den Solver lieber, Strom zu kaufen, "
        "als Komfort zu opfern. Faktor 10 Sicherheitsabstand.",
        Inches(0.6), Inches(6.1), SLIDE_W - Inches(1.2), Inches(0.85),
        bg_color=C_LIGHT_AMBER_BG, border_color=C_AMBER, fontsize=14,
    )

    add_footer(slide, page)


def slide_term_aging(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Term 4:  Batterie-Alterungskosten",
                  "Bremst übermäßiges Cycling")

    add_equation_image(
        slide,
        r"K^{\text{age}} = \frac{c^{\text{age}}}{2} \sum_{t \in T} \Delta t \cdot "
        r"\left( P^{\text{batt,ch}}_t + P^{\text{batt,dis}}_t \right)",
        Inches(2.0), Inches(1.3), target_height_in=0.9, fontsize=24,
    )

    add_equation_image(
        slide,
        r"c^{\text{age}} = \frac{C^{\text{Ersatz}} - R^{\text{EOL}}}"
        r"{N_{\text{EFC}} \cdot E^{\text{nutzbar}} \cdot \eta_{\text{rt}}}",
        Inches(3.5), Inches(2.5), target_height_in=0.85, fontsize=20,
    )

    items = [
        "C^Ersatz: Wiederbeschaffungswert (typisch 500 EUR/kWh)",
        "N_EFC: Zyklen bis EOL (LFP 6000–10000, NMC 3000–5000)",
        "η_rt: Roundtrip-Wirkungsgrad (typisch 0.85–0.95)",
        "Halbiert in der Zielfunktion: ein Vollzyklus = 1× laden + 1× entladen",
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.3,
                 top=Inches(3.9))

    add_callout_box(
        slide,
        "Wirtschaftliche Konsequenz: ein Lade-Entlade-Zyklus muss eine "
        "Preisdifferenz von mindestens c^age / η_rt überschreiten, sonst "
        "lohnt er sich nicht. Bei Defaults: ~7 ct/kWh Hürde.",
        Inches(0.6), Inches(6.1), SLIDE_W - Inches(1.2), Inches(0.85),
        bg_color=C_LIGHT_GREEN_BG, border_color=C_GREEN, fontsize=14,
    )

    add_footer(slide, page)


def slide_objective_complete(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Die vollständige Zielfunktion",
                  "Alle vier Terme zusammen")

    add_equation_image(
        slide,
        r"\min \sum_{t \in T} \Delta t \cdot "
        r"[ \pi_t P^{\text{buy}}_t - \pi^{\text{feed}} P^{\text{sell}}_t "
        r"+ c^{\text{slack}} \left( s^{\text{Heiz}}_t + s^{\text{WW}}_t \right) "
        r"+ \frac{c^{\text{age}}}{2} \left( P^{\text{batt,ch}}_t + P^{\text{batt,dis}}_t \right) "
        r"]",
        Inches(0.5), Inches(1.5), target_height_in=1.1, fontsize=20,
    )

    items = [
        "Lineare Funktion der Entscheidungsvariablen → MILP-kompatibel",
        "Pro Zeitschritt: vier kleine Beiträge, gewichtet mit Δt",
        "Summe über alle 96 Zeitschritte (24 h × 15 min) = Tageskosten",
        "Was nicht drin ist (bewusst): CO₂, WP-Verschleiß, "
        "Investitionskosten – das sind keine durch Steuerung beeinflussbaren "
        "Größen",
    ]
    add_bullets(slide, items, fontsize=18, line_spacing=1.35,
                 top=Inches(3.5))

    add_footer(slide, page)


# ============ Komponenten-Beiträge ============

def slide_section_components(prs, page):
    slide = add_blank_slide(prs)
    add_section_title(slide, "Wer trägt was bei?",
                      kicker="Teil 2 – Komponenten")


def slide_components_overview(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Komponenten und ihre MILP-Beiträge",
                  "Variablen, Constraints, Rolle in der Bilanz")

    rows = [
        ("Komponente", "Variablen", "Hauptbeitrag", "Rolle"),
        ("Netz", "P^buy, P^sell + Binär", "Disjunktion Bezug ↔ Einspeisung", "Quelle/Senke"),
        ("Batterie", "P_ch, P_dis, E + 2 Binär", "SOC-Bilanz, Lade/Entlade-Disjunktion", "Speicher"),
        ("Wärmepumpe", "P^HP, y^HP, y^start + SG1/2/3/4", "SG-Ready einziger Steuerkanal, max 8 Starts/Tag", "Wandler"),
        ("Estrich (FBH)", "E^Floor, Q^Floor_in, Q^Floor→Raum", "Bilanz mit explizitem Wärmestrom an Raum", "Speicher"),
        ("Gebäude (Raum)", "T_innen, s_low, s_high", "Raumluftbilanz + Komfortband (Mai 2026)", "Speicher"),
        ("Pufferspeicher (WW)", "E^WW, Q^WW_in + Legionellen-Bin.", "Zwei-Zonen-Verluste, SG3/SG4-Boost", "Speicher"),
        ("Wallbox", "P^WB, y^WB, SOC^EV", "Ziel-SOC zur Abfahrt, 5%/h Fahrverbrauch", "Senke"),
        ("PV / FWS / E-Auto", "(keine eigenen)", "Lieferung von Eingangs-Zeitreihen", "passiv"),
    ]

    # Tabelle bauen
    top = Inches(1.2)
    left = Inches(0.5)
    width = SLIDE_W - Inches(1.0)
    height = Inches(5.5)
    cols = 4
    rows_n = len(rows)

    table_shape = slide.shapes.add_table(rows_n, cols, left, top, width, height)
    table = table_shape.table

    col_widths = [Inches(2.5), Inches(3.5), Inches(4.6), Inches(1.7)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.name = "Calibri" if ri > 0 else "Calibri"
            r.font.size = Pt(13) if ri == 0 else Pt(12)
            r.font.bold = (ri == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_DARK_GREY
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = C_DARK_BLUE
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = C_LIGHT_GREY
            else:
                cell.fill.fore_color.rgb = C_WHITE

    add_footer(slide, page)


def slide_grid_balance(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Die Knotenbilanz Strom",
                  "Die zentrale Kopplung aller Komponenten")

    add_equation_image(
        slide,
        r"P^{\text{PV}}_t + P^{\text{buy}}_t + P^{\text{batt,dis}}_t \;=\; "
        r"P^{\text{Last}}_t + P^{\text{sell}}_t + P^{\text{batt,ch}}_t + "
        r"P^{\text{HP}}_t + \sum_w P^{\text{WB},w}_t",
        Inches(0.5), Inches(1.3), target_height_in=1.0, fontsize=18,
    )

    items = [
        {"text": "Linke Seite – Erzeugung: PV + Netzbezug + Batterie-Entladung", "bold_lead": False},
        {"text": "Rechte Seite – Verbrauch: Haushaltslast + Netzeinspeisung + "
                  "Batterie-Ladung + Wärmepumpe + Summe aller Wallboxen", "bold_lead": False},
        {"text": "Kirchhoff am AC-Knoten – muss in jedem Zeitschritt exakt erfüllt sein", "bold_lead": False},
        {"text": "Hier kommt die Sektorenkopplung Strom ↔ Wärme zustande: "
                  "PV-Überschuss kann in die WP fließen und dort als Wärme im Estrich landen", "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.35,
                 top=Inches(2.8))

    add_callout_box(
        slide,
        "Diese eine Constraint ist die zentrale Stelle, an der alle Komponenten "
        "miteinander reden. Eine neue Komponente integriert sich in EMOS Light "
        "primär über diese Bilanz – sie liefert ihren Beitrag als Term links "
        "oder rechts.",
        Inches(0.6), Inches(5.5), SLIDE_W - Inches(1.2), Inches(1.4),
        bg_color=C_LIGHT_GREEN_BG, border_color=C_GREEN, fontsize=14,
    )

    add_footer(slide, page)


def slide_battery_constraints(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Batterie – Constraints im Detail",
                  "SOC-Bilanz und Lade/Entlade-Disjunktion")

    # Constraint 1
    add_equation_image(
        slide,
        r"y^{\text{ch}}_t + y^{\text{dis}}_t \;\leq\; 1",
        Inches(0.6), Inches(1.4), target_height_in=0.55, fontsize=18,
    )
    # Constraint 2 + 3
    add_equation_image(
        slide,
        r"P^{\text{batt,ch}}_t \leq \bar{P}^{\text{ch}} \cdot y^{\text{ch}}_t, \quad "
        r"P^{\text{batt,dis}}_t \leq \bar{P}^{\text{dis}} \cdot y^{\text{dis}}_t",
        Inches(0.6), Inches(2.1), target_height_in=0.55, fontsize=18,
    )
    # SOC-Bilanz
    add_equation_image(
        slide,
        r"E^{\text{batt}}_t \;=\; E^{\text{batt}}_{t-1} "
        r"+ \eta_{\text{ch}} \cdot P^{\text{batt,ch}}_t \cdot \Delta t "
        r"- \frac{1}{\eta_{\text{dis}}} \cdot P^{\text{batt,dis}}_t \cdot \Delta t",
        Inches(0.6), Inches(2.9), target_height_in=0.75, fontsize=18,
    )

    items = [
        "Disjunktion (Zeile 1): nie gleichzeitig laden und entladen",
        "Big-M-Kopplung (Zeile 2): Leistung nur, wenn Binärvariable aktiv",
        "SOC-Bilanz (Zeile 3): Verluste auf der DC-Seite verbucht – "
        "1 kWh AC laden bringt η_ch kWh in den Speicher",
        "Bounds: E^min ≤ E_t ≤ E^max (nutzbares SoC-Fenster)",
    ]
    add_bullets(slide, items, fontsize=16, line_spacing=1.3,
                 top=Inches(4.3))

    add_callout_box(
        slide,
        "Ohne Disjunktion könnte der Solver beide Variablen > 0 setzen "
        "und die Roundtrip-Verluste als 'kostenlose Energievernichtung' "
        "missbrauchen – physikalisch unmöglich, mathematisch verlockend.",
        Inches(0.6), Inches(6.4), SLIDE_W - Inches(1.2), Inches(0.65),
        bg_color=C_LIGHT_AMBER_BG, border_color=C_AMBER, fontsize=13,
    )

    add_footer(slide, page)


def slide_heatpump_constraints(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Wärmepumpe – Constraints im Detail",
                  "Modulation, Hardwareschutz, SG-Ready")

    items = [
        {"text": "Modulationsbereich: P^HP_min · y^HP_t ≤ P^HP_t ≤ P^HP_max · y^HP_t",
         "bold_lead": True},
        {"text": "Mindestlaufzeit / Mindestpausenzeit (Hardware-Schutz)",
         "bold_lead": True},
        {"text": "Tageslimit Einschaltvorgänge (Mai 2026): "
                 "y^start_t ≥ y^HP_t − y^HP_{t-1};  Σ_{t∈Tag} y^start_t ≤ 8",
         "bold_lead": True},
        {"text": "Thermische Kopplung: Q^Floor_in_t = COP^heiz_t · P^HP,Floor_t "
                 "und Q^WW_in_t = COP^ww_t · P^HP,WW_t  (COP vorberechnet → linear)",
         "bold_lead": True},
        {"text": "SG-Ready (BWP v1.1) als einziger Steuerkanal: "
                 "Σ_i y^SGi_t = 1   und   y^HP_t + y^SG1_t = 1",
         "bold_lead": True},
        {"text": "SG3: WW-Sollwert +ΔT³. SG4: zusätzlich auch Estrich-Sollwert "
                 "+ΔT⁴ (Zwangseinschaltung mit Pufferspeicher-Boost)",
         "bold_lead": True},
    ]
    add_bullets(slide, items, fontsize=16, line_spacing=1.35,
                 top=Inches(1.3))

    add_callout_box(
        slide,
        "Trick: der COP wird VOR der Optimierung als Zeitreihe berechnet "
        "(2D-Interpolation aus aroTHERM-Kennfeld, abhängig von Außen- und "
        "Vorlauftemperatur). Damit bleibt die Beziehung zwischen elektrischer "
        "und thermischer Leistung linear – sonst wäre das Problem bilinear.",
        Inches(0.6), Inches(5.7), SLIDE_W - Inches(1.2), Inches(1.3),
        bg_color=C_LIGHT_GREEN_BG, border_color=C_GREEN, fontsize=13,
    )

    add_footer(slide, page)


def slide_thermal_storage(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Thermische Speicher (Mai 2026)",
                  "Estrich → Raum → Außen als explizite Energiekette")

    # Estrich-Bilanz mit q_floor_to_room als Variable
    add_equation_image(
        slide,
        r"E^{\text{Floor}}_t \;=\; E^{\text{Floor}}_{t-1} "
        r"+ \left( \dot{Q}^{\text{Floor,in}}_t - \dot{Q}^{\text{Floor}\to\text{Raum}}_t \right) \Delta t",
        Inches(0.6), Inches(1.25), target_height_in=0.6, fontsize=15,
    )
    # Raumluftbilanz (NEU Mai 2026)
    add_equation_image(
        slide,
        r"C^{\text{room}} (T^{\text{innen}}_t - T^{\text{innen}}_{t-1}) "
        r"= ( \dot{Q}^{\text{Floor}\to\text{Raum}}_t - \dot{Q}^{\text{Verlust}}_t ) \Delta t",
        Inches(0.6), Inches(1.95), target_height_in=0.6, fontsize=15,
    )
    # WW-Bilanz
    add_equation_image(
        slide,
        r"E^{\text{WW}}_t \;=\; E^{\text{WW}}_{t-1} "
        r"+ \left( \dot{Q}^{\text{WW,in}}_t - \dot{Q}^{\text{WW,bedarf}}_t \right) \Delta t "
        r"- \dot{Q}^{\text{fix}} \Delta t - \mu^{\text{rel}} E^{\text{WW}}_{t-1} \Delta t",
        Inches(0.6), Inches(2.65), target_height_in=0.7, fontsize=13,
    )

    items = [
        "Estrich → Raum: Q^Floor→Raum eigene MILP-Variable; Bilanz affin in T_floor und T_innen",
        "Raumluft: T_innen ist seit Mai 2026 eine MILP-Zustandsvariable",
        "Komfortband: T_innen ∈ [T_min, T_max] als Soft-Constraint, Slack-Penalty 500 ct/kWh",
        "WW-Speicher: Zwei-Zonen mit geometriebasierten Verlusten (unverändert)",
        "Fallback ohne Building: alte Verlustraten-Bilanz mit λ · E (rückwärtskompatibel)",
    ]
    add_bullets(slide, items, fontsize=13, line_spacing=1.25,
                 top=Inches(3.6))

    add_footer(slide, page)


def slide_wallbox(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Wallbox – Constraints pro Fahrzeug",
                  "Modulation, Anwesenheit, Mindestlademenge")

    items = [
        "Modulationsbereich: P^WB,min · y^WB_t ≤ P^WB_t ≤ P^WB,max · y^WB_t",
        "EV-Anwesenheit: P^WB_t = 0, wenn das Fahrzeug nicht am Stecker hängt",
        {"text": "Mindestlademenge: garantiert die geforderte Reichweite", "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=18, line_spacing=1.35, top=Inches(1.4))

    add_equation_image(
        slide,
        r"\sum_{t \in T} P^{\text{WB},w}_t \cdot \Delta t \;\geq\; "
        r"\frac{\Delta\text{SOC} \cdot E^{\text{EV,kap}}}{\eta_{\text{WB}}}",
        Inches(2.0), Inches(3.6), target_height_in=0.95, fontsize=22,
    )

    items2 = [
        "ΔSOC = SOC_ziel − SOC_aktuell, abgeleitet aus Mindestreichweite × Verbrauch",
        "Pro Wallbox eigener Block – mehrere Wallboxen koexistieren als unabhängige Blöcke",
        "Optional: Strompreis-Perzentil-Filter (nur in günstigsten X % der "
        "Anwesenheitsstunden laden)",
    ]
    add_bullets(slide, items2, fontsize=16, line_spacing=1.3, top=Inches(5.0))

    add_footer(slide, page)


# ============ Solver ============

def slide_section_solver(prs, page):
    slide = add_blank_slide(prs)
    add_section_title(slide, "Wie löst der Solver das?",
                      kicker="Teil 3 – Algorithmus")


def slide_milp_class(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Die Problemklasse: MILP",
                  "Mixed-Integer Linear Programming")

    add_equation_image(
        slide,
        r"\min_x \; c^\top x \quad \text{s.t.} \quad A x \leq b, \quad "
        r"l \leq x \leq u, \quad x_j \in \mathbb{Z} \; \text{für} \; j \in J",
        Inches(1.0), Inches(1.4), target_height_in=0.9, fontsize=20,
    )

    items = [
        "Linear: Zielfunktion und alle Constraints sind lineare Ausdrücke",
        "Konvex: ohne Ganzzahligkeit wäre der zulässige Bereich ein "
        "konvexer Polyeder – Optimum in einer Ecke (Fundamentalsatz LP)",
        "Ganzzahlig: Binärvariablen für Schalt-Entscheidungen "
        "(WP an/aus, Batterie laden/entladen, SG-Ready 1/3)",
        "Ohne Ganzzahligkeit polynomial lösbar – mit ist es NP-schwer "
        "(theoretisch). In der Praxis: HiGHS löst unser Modell in &lt; 1 s",
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.35,
                 top=Inches(3.0))

    add_footer(slide, page)


def slide_lp_relaxation(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Schritt 1:  LP-Relaxation",
                  "Untere Schranke durch Lockerung der Ganzzahligkeit")

    items = [
        "Binärvariablen y ∈ {0,1} werden zu kontinuierlichen y ∈ [0,1] gemacht",
        "Resultierendes Problem: reines LP – polynomial lösbar mit Simplex",
        "LP-Lösung typisch fraktional (z.B. y^HP_14:00 = 0.42) – "
        "physikalisch unsinnig, aber liefert eine untere Schranke",
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.35,
                 top=Inches(1.3))

    add_equation_image(
        slide,
        r"f^*_{\text{LP}} \;\leq\; f^*_{\text{MILP}}",
        Inches(4.5), Inches(3.5), target_height_in=0.8, fontsize=24,
    )

    add_callout_box(
        slide,
        "Die LP-Relaxation ist der Kern der Beweisführung: sie garantiert, "
        "dass die echte MILP-Lösung nicht günstiger sein kann als der "
        "LP-Optimum-Wert. Damit kann der Solver gezielt suchen, statt "
        "alle 10^174 möglichen Binärbelegungen aufzuzählen.",
        Inches(0.6), Inches(5.0), SLIDE_W - Inches(1.2), Inches(1.6),
        bg_color=C_LIGHT_BLUE_BG, border_color=C_MID_BLUE, fontsize=14,
    )

    add_footer(slide, page)


def slide_branch_and_bound(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Schritt 2:  Branch-and-Bound",
                  "Verzweigen und kluges Abschneiden")

    items = [
        {"text": "Verzweigen: wähle eine fraktionale Binärvariable und teile "
                 "das Problem in zwei Unterprobleme (y=0 und y=1)", "bold_lead": False},
        {"text": "Jedes Unterproblem ist wieder ein LP – kann gelöst werden", "bold_lead": False},
        {"text": "Pruning: Subbäume verwerfen, sobald deren LP-Bound die "
                 "beste bisherige ganzzahlige Lösung (Inkumbente) nicht "
                 "mehr unterbieten kann", "bold_lead": False},
        {"text": "Bei EMOS Light: statt 10^174 Möglichkeiten werden "
                 "tatsächlich nur einige hundert Knoten besucht", "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.35, top=Inches(1.3))

    add_callout_box(
        slide,
        "Anschauung: wie ein Detektiv. Statt jeden Verdächtigen einzeln zu "
        "befragen, schließt der Solver ganze Personengruppen aus, weil "
        "deren optimistische Schätzung schon teurer ist als das beste "
        "bekannte echte Ergebnis.",
        Inches(0.6), Inches(5.4), SLIDE_W - Inches(1.2), Inches(1.5),
        bg_color=C_LIGHT_GREEN_BG, border_color=C_GREEN, fontsize=14,
    )

    add_footer(slide, page)


def slide_convergence(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Konvergenz:  zwei Schranken treffen sich",
                  "MIP-Gap → 0 = Optimum bewiesen")

    add_equation_image(
        slide,
        r"f_{\text{LB}} \;\leq\; f^* \;\leq\; f_{\text{UB}}, \qquad "
        r"\text{Gap} = \frac{f_{\text{UB}} - f_{\text{LB}}}{|f_{\text{UB}}|}",
        Inches(1.0), Inches(1.3), target_height_in=0.9, fontsize=20,
    )

    items = [
        {"text": "Untere Schranke (LB): bestes LP-Ergebnis offener Knoten – "
                 "steigt durch Pruning monoton", "bold_lead": False},
        {"text": "Obere Schranke (UB): bester gefundener echter Plan (Inkumbente) – "
                 "sinkt monoton durch Heuristiken (Diving, Feasibility Pump)", "bold_lead": False},
        {"text": "Wenn UB = LB (oder Gap < Toleranz): Optimum bewiesen", "bold_lead": False},
        {"text": "Bei EMOS Light: typisch &lt; 1 s bis Gap &lt; 0.01 %", "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.35, top=Inches(2.8))

    add_callout_box(
        slide,
        "MILP liefert eine BEWIESEN optimale Lösung – nicht 'wahrscheinlich gut' "
        "wie bei Heuristiken oder neuronalen Netzen. Das ist der entscheidende "
        "Vorteil gegenüber RL- oder ML-basierten Ansätzen.",
        Inches(0.6), Inches(5.7), SLIDE_W - Inches(1.2), Inches(1.3),
        bg_color=C_LIGHT_BLUE_BG, border_color=C_MID_BLUE, fontsize=14,
    )

    add_footer(slide, page)


def slide_lifecycle(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Lebenslauf eines Optimierungslaufs",
                  "Was passiert in unter einer Sekunde")

    rows = [
        ("t [ms]", "Phase", "Aktion"),
        ("0",     "Modellbau",       "PuLP erzeugt ~1.700 Variablen und Constraints"),
        ("100",   "I/O",             "LP-Datei schreiben, HiGHS starten"),
        ("130",   "Presolve",        "Modell um 30–60 % schrumpfen"),
        ("180",   "Wurzel-LP",       "Erste LP-Relaxation, LB ≈ 4,21 EUR"),
        ("200",   "Schnittebenen",   "~30 Cuts hinzufügen, LB steigt"),
        ("250",   "Heuristik",       "Diving findet erste Inkumbente, UB ≈ 4,71 EUR"),
        ("280",   "Branch-and-Bound", "Verzweigen, prunen, Heuristiken"),
        ("680",   "Konvergenz",      "Gap < 0.01 %, Optimum bewiesen"),
        ("700",   "Rücklesen",       "Variablenwerte → OptimizationResult"),
    ]

    top = Inches(1.2)
    left = Inches(0.5)
    width = SLIDE_W - Inches(1.0)
    height = Inches(5.7)
    cols = 3
    rows_n = len(rows)

    table_shape = slide.shapes.add_table(rows_n, cols, left, top, width, height)
    table = table_shape.table

    col_widths = [Inches(1.5), Inches(2.7), Inches(8.1)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = val
            r.font.name = "Calibri"
            r.font.size = Pt(14) if ri == 0 else Pt(13)
            r.font.bold = (ri == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_DARK_GREY
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = C_DARK_BLUE
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = C_LIGHT_GREY
            else:
                cell.fill.fore_color.rgb = C_WHITE

    add_footer(slide, page)


def slide_modelsize(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Modellgröße bei EMOS Light",
                  "Was der Solver tatsächlich vor sich hat")

    rows = [
        ("Konfiguration", "Variablen", "davon binär", "Constraints", "Lösungszeit"),
        ("Nur Netz + Last (24 h, 15 min)", "≈ 290", "≈ 100", "≈ 200", "~ 50 ms"),
        ("+ PV + Batterie", "≈ 770", "≈ 290", "≈ 670", "~ 200 ms"),
        ("+ WP + FBH + WW + SG-Ready", "≈ 1700", "≈ 580", "≈ 1700", "~ 500 ms"),
        ("+ 2 Wallboxen (Vollszenario)", "≈ 2100", "≈ 770", "≈ 2200", "~ 700 ms"),
        ("MPC-Mehrtages (35 h Horizont)", "≈ 2400", "≈ 850", "≈ 2500", "~ 1500 ms"),
    ]

    top = Inches(1.2)
    left = Inches(0.5)
    width = SLIDE_W - Inches(1.0)
    height = Inches(4.5)
    cols = 5
    rows_n = len(rows)

    table_shape = slide.shapes.add_table(rows_n, cols, left, top, width, height)
    table = table_shape.table

    col_widths = [Inches(4.5), Inches(1.8), Inches(1.8), Inches(2.0), Inches(2.2)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    for ri, row_data in enumerate(rows):
        for ci, val in enumerate(row_data):
            cell = table.cell(ri, ci)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.name = "Calibri"
            r.font.size = Pt(13)
            r.font.bold = (ri == 0)
            r.font.color.rgb = C_WHITE if ri == 0 else C_DARK_GREY
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = C_DARK_BLUE
            elif ri % 2 == 0:
                cell.fill.fore_color.rgb = C_LIGHT_GREY
            else:
                cell.fill.fore_color.rgb = C_WHITE

    add_callout_box(
        slide,
        "Selbst das volle Szenario lässt sich im Sekundenbereich lösen. "
        "Das ist der Grund, warum MPC praktikabel ist – wir können stündlich "
        "neu rechnen, ohne dass das Hardware-Probleme verursacht.",
        Inches(0.6), Inches(5.9), SLIDE_W - Inches(1.2), Inches(1.0),
        bg_color=C_LIGHT_GREEN_BG, border_color=C_GREEN, fontsize=14,
    )

    add_footer(slide, page)


# ============ Schluss ============

def slide_section_takeaway(prs, page):
    slide = add_blank_slide(prs)
    add_section_title(slide, "Was steht – und was kommt",
                      kicker="Teil 4")


def slide_status(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Stand der Implementierung",
                  "Was funktioniert, was gerade refactored wurde")

    items = [
        {"text": "Alle 9 Komponenten implementiert – mit realen Parametern "
                 "(aroTHERM-Kennfeld, BWP-SG-Ready, DIN-basierte Gebäudewerte)",
         "bold_lead": False},
        {"text": "MILP-Optimierer + MPC-Schleife + Baseline funktionieren",
         "bold_lead": False},
        {"text": "Dashboard (Streamlit) mit allen Konfigurationsfeldern, "
                 "Plots, YAML-Import/-Export",
         "bold_lead": False},
        {"text": "Refactoring abgeschlossen: zwei-stufige Basisklasse "
                 "(Component vs MILPComponent), MILP-Helfer-Modul mit 8 "
                 "wiederverwendbaren Bausteinen",
         "bold_lead": False},
        {"text": "Effekt: Komponenten-Code 26 % kürzer, eine neue Komponente "
                 "lässt sich in 2–4 h statt einem Tag implementieren",
         "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.4, top=Inches(1.4))

    add_callout_box(
        slide,
        "Praktischer Effekt des Refactorings: die häufig duplizierten Muster "
        "(z.B. SOC-Bilanz mit t=0-Sonderfall, EIN/AUS-Power-Link) stehen "
        "einmal zentral. Komponenten schreiben Intent statt Boilerplate.",
        Inches(0.6), Inches(5.7), SLIDE_W - Inches(1.2), Inches(1.3),
        fontsize=14,
    )

    add_footer(slide, page)


def slide_outlook(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Ausblick",
                  "Die nächsten Schritte")

    items = [
        {"text": "2R2C-Gebäudemodell: separater Raum-/Wand-Knoten mit "
                 "Transmissions- und Lüftungsverlusten als Constraints",
         "bold_lead": False},
        {"text": "Hybrid-PV-Prognose: Day-Ahead-Wetter + HTW-Intraday-"
                 "Korrektur – soll typische Fehler bei Bewölkung halbieren",
         "bold_lead": False},
        {"text": "Adaptiver MPC-Horizont: vormittags nur Resttag, ab 13 Uhr "
                 "(wenn EPEX-Preise für morgen da sind) voller 35-h-Horizont",
         "bold_lead": False},
        {"text": "Optimizer-Refactoring Stufe 2: generische Komponenten-"
                 "Schleife, dict-basiertes Result, Plugin-Registry",
         "bold_lead": False},
        {"text": "Erweiterte Komponenten auf der Wunschliste: Heizstab "
                 "(~½ Tag), Solarthermie (~1 Tag), BHKW (~2–3 Tage)",
         "bold_lead": False},
        {"text": "Validierung gegen Realdaten aus dem Pilotgebäude (InfluxDB)",
         "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=16, line_spacing=1.35, top=Inches(1.3))

    add_footer(slide, page)


def slide_takehome(prs, page):
    slide = add_blank_slide(prs)
    add_title_bar(slide, "Take-Home Messages",
                  "Was du mitnehmen solltest")

    items = [
        {"text": "MILP, nicht KI: bewiesen optimale Lösungen, erklärbar, "
                 "ohne Trainingsdaten – Standard in der industriellen "
                 "Energieoptimierung",
         "bold_lead": False},
        {"text": "Zielfunktion = vier lineare Terme: Netzkosten, Einspeise-"
                 "Erlös, Komfort-Strafe, Batterie-Alterung",
         "bold_lead": False},
        {"text": "Knotenbilanz als zentrale Kopplung: alle Komponenten "
                 "treffen sich in einer Gleichung pro Zeitschritt",
         "bold_lead": False},
        {"text": "Linearität ist kein Zufall, sondern Modellentscheidung: "
                 "COP vorberechnet, Big-M-Disjunktionen, Slack mit Penalty",
         "bold_lead": False},
        {"text": "HiGHS löst typisch in unter einer Sekunde – damit MPC "
                 "praktikabel und für Echtbetrieb geeignet",
         "bold_lead": False},
        {"text": "Modulare Architektur nach Refactoring: neue Komponenten "
                 "in Stunden statt Tagen integrierbar",
         "bold_lead": False},
    ]
    add_bullets(slide, items, fontsize=17, line_spacing=1.4, top=Inches(1.3))

    add_footer(slide, page)


def slide_qa(prs, page):
    slide = add_blank_slide(prs)

    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK_BLUE
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.8), SLIDE_W - Inches(2.0), Inches(2.0),
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Fragen?"
    r.font.size = Pt(80)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    r.font.name = "Calibri"

    tb2 = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.5), SLIDE_W - Inches(2.0), Inches(1.5),
    )
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "Diskussion und Vertiefung"
    r2.font.size = Pt(24)
    r2.font.color.rgb = RGBColor(0xCC, 0xD7, 0xEC)
    r2.font.name = "Calibri"


# ----------------------------------------------------------------------
# Document assembly
# ----------------------------------------------------------------------

def build_pptx(out_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Folien
    slide_cover(prs)  # 1
    slide_agenda(prs, 2)  # 2
    slide_what_does_it_do(prs, 3)  # 3
    slide_section_obj_func(prs, 4)  # 4
    slide_objective_overview(prs, 5)  # 5
    slide_term_grid(prs, 6)  # 6
    slide_term_slack(prs, 7)  # 7
    slide_term_aging(prs, 8)  # 8
    slide_objective_complete(prs, 9)  # 9
    slide_section_components(prs, 10)  # 10
    slide_components_overview(prs, 11)  # 11
    slide_grid_balance(prs, 12)  # 12
    slide_battery_constraints(prs, 13)  # 13
    slide_heatpump_constraints(prs, 14)  # 14
    slide_thermal_storage(prs, 15)  # 15
    slide_wallbox(prs, 16)  # 16
    slide_section_solver(prs, 17)  # 17
    slide_milp_class(prs, 18)  # 18
    slide_lp_relaxation(prs, 19)  # 19
    slide_branch_and_bound(prs, 20)  # 20
    slide_convergence(prs, 21)  # 21
    slide_lifecycle(prs, 22)  # 22
    slide_modelsize(prs, 23)  # 23
    slide_section_takeaway(prs, 24)  # 24
    slide_status(prs, 25)  # 25
    slide_outlook(prs, 26)  # 26
    slide_takehome(prs, 27)  # 27
    slide_qa(prs, 28)  # 28

    prs.save(out_path)


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "EMOS_Light_Optimierer.pptx"
    out_abs = os.path.abspath(out)
    build_pptx(out_abs)
    print(f"PowerPoint geschrieben: {out_abs}")
